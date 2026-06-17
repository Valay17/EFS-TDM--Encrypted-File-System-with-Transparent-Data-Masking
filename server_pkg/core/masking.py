"""
Regex-based masking engine.

Supported formats for masking:
  Text:        .txt, .log, .tex, .bib
  CSV:         .csv  (stdlib csv module — handles multi-line quoted fields)
  Spreadsheet: .xlsx, .xls  (openpyxl / xlrd; xls output is saved as .xlsx)
  SQL:         .sql, .dump
  PDF:         .pdf  (PyMuPDF text-layer path; OCR path via pytesseract if any
                      page is image-based)
  Word:        .docx  (.doc unsupported — convert to .docx first)
  ODF:         .odt, .ods, .odp  (odfpy)
  Presentation:.pptx  (.ppt unsupported — convert to .pptx first)

Encrypt-only (no masking):
  Images:  .png .jpg .jpeg .gif .bmp .webp .tiff .tif .heic
  Video:   .mp4 .mkv .avi .mov .wmv .flv .webm
  Audio:   .mp3 .wav .aac .flac .ogg .m4a
  Archive: .zip .tar .gz .bz2 .xz .7z .rar

Masking rules are loaded from config/masking_rules.json.
"""

import csv
import io
import json
import os
import re
import threading
from pathlib import Path

import fitz  # PyMuPDF
import pytesseract
from PIL import Image
from docx import Document
from odf import text as odf_text
from odf import teletype
from odf.opendocument import load as odf_load
from odf.table import Table, TableRow, TableCell
from openpyxl import load_workbook
from pptx import Presentation
from pytesseract import Output

RULES_PATH = Path(__file__).parent.parent / "config" / "masking_rules.json"

_ENCRYPT_ONLY_EXTS = {
    ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".tiff", ".tif", ".heic",
    ".mp4", ".mkv", ".avi", ".mov", ".wmv", ".flv", ".webm",
    ".mp3", ".wav", ".aac", ".flac", ".ogg", ".m4a",
    ".zip", ".tar", ".gz", ".bz2", ".xz", ".7z", ".rar",
}


# ---------------------------------------------------------------------------
# ReDoS-safe regex application
# ---------------------------------------------------------------------------

def _safe_sub(pattern, replacement: str, text: str, timeout_secs: int = 5) -> str:
    result = [text]
    error = [None]

    def _run():
        try:
            result[0] = pattern.sub(replacement, text)
        except Exception as e:
            error[0] = e

    t = threading.Thread(target=_run, daemon=True)
    t.start()
    t.join(timeout=timeout_secs)
    if t.is_alive():
        raise RuntimeError(f"Regex pattern timed out after {timeout_secs}s — possible ReDoS")
    if error[0]:
        raise error[0]
    return result[0]


# ---------------------------------------------------------------------------
# Rules loading and compilation
# ---------------------------------------------------------------------------

def _validate_rules(rules_config: dict) -> None:
    if not isinstance(rules_config, dict):
        raise ValueError("Rules config must be a dict")
    if "rules" not in rules_config or not isinstance(rules_config["rules"], dict):
        raise ValueError("Rules config must have a 'rules' key that is a dict")
    if "role_profiles" not in rules_config or not isinstance(rules_config["role_profiles"], dict):
        raise ValueError("Rules config must have a 'role_profiles' key that is a dict")
    for name, rule in rules_config["rules"].items():
        if not isinstance(rule.get("pattern"), str):
            raise ValueError(f"Rule '{name}' must have a 'pattern' string")
        if not isinstance(rule.get("replacement"), str):
            raise ValueError(f"Rule '{name}' must have a 'replacement' string")


def load_rules() -> dict:
    with open(RULES_PATH, "r") as f:
        rules_config = json.load(f)
    _validate_rules(rules_config)
    return rules_config


def _compile_rules(rule_names: list, rules_config: dict) -> list:
    compiled = []
    for name in rule_names:
        rule = rules_config["rules"].get(name)
        if rule:
            compiled.append((re.compile(rule["pattern"]), rule["replacement"]))
    return compiled


def _apply_rules(text: str, compiled_rules: list) -> str:
    for pattern, replacement in compiled_rules:
        text = _safe_sub(pattern, replacement, text)
    return text


def _get_rules_for_role(role: str, rules_config: dict) -> list:
    rule_names = rules_config["role_profiles"].get(role, [])
    return _compile_rules(rule_names, rules_config)


def _apply_name_masking(text: str, rules_config: dict) -> str:
    replacement = rules_config["name_masking"]["replacement"]
    return re.sub(r"\[NAME\]", replacement, text)


# ---------------------------------------------------------------------------
# Public API — text / string masking
# ---------------------------------------------------------------------------

def mask_text(text, role="Analyst", mask_names=False, rules_config=None):
    if rules_config is None:
        rules_config = load_rules()
    compiled = _get_rules_for_role(role, rules_config)
    result = _apply_rules(text, compiled)
    if mask_names:
        result = _apply_name_masking(result, rules_config)
    return result


def mask_txt(content, role="Analyst", mask_names=False, rules_config=None):
    return mask_text(content, role=role, mask_names=mask_names, rules_config=rules_config)


def mask_sql(content, role="Analyst", mask_names=False, rules_config=None):
    return mask_text(content, role=role, mask_names=mask_names, rules_config=rules_config)


# ---------------------------------------------------------------------------
# CSV masking — stdlib csv module (handles multi-line quoted fields correctly)
# ---------------------------------------------------------------------------

def mask_csv(content, role="Analyst", mask_names=False, rules_config=None):
    if rules_config is None:
        rules_config = load_rules()
    compiled = _get_rules_for_role(role, rules_config)

    reader = csv.reader(io.StringIO(content))
    out = io.StringIO()
    writer = csv.writer(out)

    for i, row in enumerate(reader):
        if i == 0:
            writer.writerow(row)
            continue
        masked_row = []
        for cell in row:
            cell = _apply_rules(cell, compiled)
            if mask_names:
                cell = _apply_name_masking(cell, rules_config)
            masked_row.append(cell)
        writer.writerow(masked_row)

    return out.getvalue()


# ---------------------------------------------------------------------------
# XLSX masking — openpyxl (preserves formatting)
# ---------------------------------------------------------------------------

def mask_xlsx(input_path, output_path, role="Analyst", mask_names=False, rules_config=None):
    if rules_config is None:
        rules_config = load_rules()
    compiled = _get_rules_for_role(role, rules_config)

    wb = load_workbook(input_path)
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                if cell.value is not None and isinstance(cell.value, str):
                    v = _apply_rules(cell.value, compiled)
                    if mask_names:
                        v = _apply_name_masking(v, rules_config)
                    cell.value = v
    wb.save(output_path)


# ---------------------------------------------------------------------------
# XLS masking — xlrd read, openpyxl write (output is .xlsx)
# ---------------------------------------------------------------------------

def mask_xls(input_path, output_path, role="Analyst", mask_names=False, rules_config=None):
    import xlrd
    from openpyxl import Workbook

    if rules_config is None:
        rules_config = load_rules()
    compiled = _get_rules_for_role(role, rules_config)

    rb = xlrd.open_workbook(input_path)
    wb = Workbook()
    wb.remove(wb.active)

    for sheet_idx in range(rb.nsheets):
        rs = rb.sheet_by_index(sheet_idx)
        ws = wb.create_sheet(title=rs.name)
        for row_idx in range(rs.nrows):
            row_vals = []
            for col_idx in range(rs.ncols):
                cell = rs.cell(row_idx, col_idx)
                val = cell.value
                if isinstance(val, str):
                    val = _apply_rules(val, compiled)
                    if mask_names:
                        val = _apply_name_masking(val, rules_config)
                row_vals.append(val)
            ws.append(row_vals)

    wb.save(output_path)


# ---------------------------------------------------------------------------
# DOCX masking — python-docx
# ---------------------------------------------------------------------------

def mask_docx(input_path, output_path, role="Analyst", mask_names=False, rules_config=None):
    if rules_config is None:
        rules_config = load_rules()
    compiled = _get_rules_for_role(role, rules_config)

    doc = Document(input_path)

    def _mask_para(para):
        full = para.text
        if not full.strip():
            return
        masked = _apply_rules(full, compiled)
        if mask_names:
            masked = _apply_name_masking(masked, rules_config)
        if masked != full and para.runs:
            # Write entire masked text into the first run, clear the rest
            para.runs[0].text = masked
            for run in para.runs[1:]:
                run.text = ""

    for para in doc.paragraphs:
        _mask_para(para)

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    _mask_para(para)

    for section in doc.sections:
        for header in (section.header, section.footer):
            if header is not None:
                for para in header.paragraphs:
                    _mask_para(para)

    doc.save(output_path)


# ---------------------------------------------------------------------------
# ODT / ODS / ODP masking — odfpy
# ---------------------------------------------------------------------------

def _mask_odf_text_elements(doc_obj, compiled, mask_names, rules_config):
    for elem in list(doc_obj.getElementsByType(odf_text.P)):
        original = teletype.extractText(elem)
        if not original.strip():
            continue
        masked = _apply_rules(original, compiled)
        if mask_names:
            masked = _apply_name_masking(masked, rules_config)
        if masked != original:
            parent = elem.parentNode
            if parent is None:
                continue
            new_p = odf_text.P()
            new_p.addText(masked)
            parent.insertBefore(new_p, elem)
            parent.removeChild(elem)


def mask_odt(input_path, output_path, role="Analyst", mask_names=False, rules_config=None):
    if rules_config is None:
        rules_config = load_rules()
    compiled = _get_rules_for_role(role, rules_config)
    doc = odf_load(input_path)
    _mask_odf_text_elements(doc, compiled, mask_names, rules_config)
    doc.save(output_path)


def mask_odp(input_path, output_path, role="Analyst", mask_names=False, rules_config=None):
    if rules_config is None:
        rules_config = load_rules()
    compiled = _get_rules_for_role(role, rules_config)
    doc = odf_load(input_path)
    _mask_odf_text_elements(doc, compiled, mask_names, rules_config)
    doc.save(output_path)


def mask_ods(input_path, output_path, role="Analyst", mask_names=False, rules_config=None):
    if rules_config is None:
        rules_config = load_rules()
    compiled = _get_rules_for_role(role, rules_config)
    doc = odf_load(input_path)

    for cell in doc.getElementsByType(TableCell):
        for p in cell.getElementsByType(odf_text.P):
            original = teletype.extractText(p)
            if not original.strip():
                continue
            masked = _apply_rules(original, compiled)
            if mask_names:
                masked = _apply_name_masking(masked, rules_config)
            if masked != original:
                parent = p.parentNode
                if parent is None:
                    continue
                new_p = odf_text.P()
                new_p.addText(masked)
                parent.insertBefore(new_p, p)
                parent.removeChild(p)

    doc.save(output_path)


# ---------------------------------------------------------------------------
# PPTX masking — python-pptx
# ---------------------------------------------------------------------------

def mask_pptx(input_path, output_path, role="Analyst", mask_names=False, rules_config=None):
    if rules_config is None:
        rules_config = load_rules()
    compiled = _get_rules_for_role(role, rules_config)

    prs = Presentation(input_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                full = "".join(run.text for run in para.runs)
                if not full.strip():
                    continue
                masked = _apply_rules(full, compiled)
                if mask_names:
                    masked = _apply_name_masking(masked, rules_config)
                if masked != full and para.runs:
                    para.runs[0].text = masked
                    for run in para.runs[1:]:
                        run.text = ""

    prs.save(output_path)


# ---------------------------------------------------------------------------
# PDF masking — hybrid text / OCR path
# ---------------------------------------------------------------------------

def _is_image_based_page(page) -> bool:
    text = page.get_text().strip()
    images = page.get_images()
    return len(text) < 50 and len(images) > 0


def _mask_pdf_text_path(doc, patterns, output_path):
    for page in doc:
        for widget in page.widgets():
            if widget.field_value:
                original = str(widget.field_value)
                masked = original
                for pattern in patterns:
                    masked = pattern.sub("*" * len(original), masked)
                if masked != original:
                    widget.field_value = masked
                    widget.update()

        page_text = page.get_text("text")
        for pattern in patterns:
            for match in pattern.finditer(page_text):
                matched_str = match.group()
                for rect in page.search_for(matched_str):
                    page.add_redact_annot(rect, fill=(0, 0, 0))

        page.apply_redactions()

    doc.save(output_path, garbage=4, deflate=True)
    doc.close()


def _mask_pdf_ocr_path(doc, patterns, output_path):
    masked_images = []

    for page in doc:
        pix = page.get_pixmap(dpi=300)
        img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)

        data = pytesseract.image_to_data(img, output_type=Output.DICT)
        from PIL import ImageDraw
        draw = ImageDraw.Draw(img)

        n_boxes = len(data["text"])
        for i in range(n_boxes):
            word = data["text"][i]
            if not word.strip():
                continue
            for pattern in patterns:
                if pattern.search(word):
                    x = data["left"][i]
                    y = data["top"][i]
                    w = data["width"][i]
                    h = data["height"][i]
                    draw.rectangle([x, y, x + w, y + h], fill=(0, 0, 0))
                    break

        masked_images.append(img)

    doc.close()

    if masked_images:
        masked_images[0].save(
            output_path,
            save_all=True,
            append_images=masked_images[1:],
            format="PDF",
        )


def mask_pdf(input_path, output_path, role="Analyst", mask_names=False, rules_config=None):
    if not os.path.exists(input_path):
        raise FileNotFoundError(f"Input PDF not found: {input_path}")

    if rules_config is None:
        rules_config = load_rules()

    rule_names = rules_config["role_profiles"].get(role, [])
    patterns = []
    for name in rule_names:
        rule = rules_config["rules"].get(name)
        if rule:
            patterns.append(re.compile(rule["pattern"]))

    doc = fitz.open(input_path)

    use_ocr = any(_is_image_based_page(page) for page in doc)

    if use_ocr:
        _mask_pdf_ocr_path(doc, patterns, output_path)
    else:
        _mask_pdf_text_path(doc, patterns, output_path)


# ---------------------------------------------------------------------------
# mask_file — main dispatch
# ---------------------------------------------------------------------------

def mask_file(input_path, role="Analyst", mask_names=False, output_path=None, rules_config=None):
    """
    Dispatch masking to the correct handler based on file extension.

    For text-based files (csv, txt, sql, tex, bib):
        Returns masked string. If output_path given, also writes to disk.

    For binary document files (pdf, xlsx, xls, docx, odt, ods, odp, pptx):
        output_path is required. Returns None (writes directly to output_path).

    Raises:
        FileNotFoundError: input_path does not exist.
        ValueError:        Unsupported or encrypt-only format, or missing output_path.
    """
    if not os.path.exists(input_path):
        raise FileNotFoundError(f"File not found: {input_path}")

    if rules_config is None:
        rules_config = load_rules()

    ext = Path(input_path).suffix.lower()

    # Encrypt-only formats — masking not supported
    if ext in _ENCRYPT_ONLY_EXTS:
        raise ValueError(
            f"Format {ext} is encrypt-only; masking is not supported for "
            "binary/media/archive files. Use encrypt/send to store securely."
        )

    # Explicitly unsupported formats
    if ext == ".doc":
        raise ValueError("Unsupported format: .doc — convert to .docx first")
    if ext == ".ppt":
        raise ValueError("Unsupported format: .ppt — convert to .pptx first")
    if ext == ".rtf":
        raise ValueError("Unsupported format: .rtf — convert to .docx or .txt first")

    # Binary document formats — require output_path
    binary_formats = {".pdf", ".xlsx", ".xls", ".docx", ".odt", ".ods", ".odp", ".pptx"}
    if ext in binary_formats and not output_path:
        raise ValueError(f"output_path is required for {ext} masking")

    if ext == ".pdf":
        mask_pdf(input_path, output_path, role=role, mask_names=mask_names, rules_config=rules_config)
        return None

    if ext == ".xlsx":
        mask_xlsx(input_path, output_path, role=role, mask_names=mask_names, rules_config=rules_config)
        return None

    if ext == ".xls":
        # Output is always .xlsx — rename output_path extension if needed
        out = output_path
        if Path(out).suffix.lower() == ".xls":
            out = str(Path(out).with_suffix(".xlsx"))
        mask_xls(input_path, out, role=role, mask_names=mask_names, rules_config=rules_config)
        return None

    if ext == ".docx":
        mask_docx(input_path, output_path, role=role, mask_names=mask_names, rules_config=rules_config)
        return None

    if ext == ".odt":
        mask_odt(input_path, output_path, role=role, mask_names=mask_names, rules_config=rules_config)
        return None

    if ext == ".ods":
        mask_ods(input_path, output_path, role=role, mask_names=mask_names, rules_config=rules_config)
        return None

    if ext == ".odp":
        mask_odp(input_path, output_path, role=role, mask_names=mask_names, rules_config=rules_config)
        return None

    if ext == ".pptx":
        mask_pptx(input_path, output_path, role=role, mask_names=mask_names, rules_config=rules_config)
        return None

    # Text-based formats
    with open(input_path, "r", encoding="utf-8", errors="replace") as f:
        content = f.read()

    if ext == ".csv":
        try:
            masked = mask_csv(content, role=role, mask_names=mask_names, rules_config=rules_config)
        except Exception:
            # Malformed CSV (bad quoting, NUL bytes, etc.) — fall back to plain-text masking
            masked = mask_txt(content, role=role, mask_names=mask_names, rules_config=rules_config)
    elif ext in (".txt", ".log", ".tex", ".bib"):
        masked = mask_txt(content, role=role, mask_names=mask_names, rules_config=rules_config)
    elif ext in (".sql", ".dump"):
        masked = mask_sql(content, role=role, mask_names=mask_names, rules_config=rules_config)
    else:
        # Unknown / no extension — treat as plain text
        masked = mask_txt(content, role=role, mask_names=mask_names, rules_config=rules_config)

    if output_path:
        os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(masked)

    return masked
