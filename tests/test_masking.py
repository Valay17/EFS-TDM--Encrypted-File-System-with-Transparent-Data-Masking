"""
Unit tests for core/masking.py.
Covers: CSV, TXT, SQL, each masking rule, admin passthrough, name flag,
XLSX, DOCX, ODT, ODS, ODP, PPTX, PDF hybrid, encrypt-only guard,
unsupported format guard, LaTeX/BIB plain-text path.
"""

import io
import os
import sys
import tempfile
import unittest

sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

from core.masking import (
    load_rules,
    mask_text,
    mask_csv,
    mask_txt,
    mask_sql,
    mask_file,
)

import fitz


RULES = load_rules()


class TestMaskText(unittest.TestCase):

    def test_ssn_masked(self):
        result = mask_text("SSN: 123-45-6789", role="Analyst", rules_config=RULES)
        self.assertNotIn("123-45-6789", result)
        self.assertIn("***-**-****", result)

    def test_email_masked(self):
        result = mask_text("Email: john.doe@example.com", role="Analyst", rules_config=RULES)
        self.assertNotIn("john.doe@example.com", result)
        self.assertIn("***@***.***", result)

    def test_credit_card_masked(self):
        result = mask_text("Card: 4111-1111-1111-1111", role="Analyst", rules_config=RULES)
        self.assertNotIn("4111-1111-1111-1111", result)
        self.assertIn("****-****-****-****", result)

    def test_credit_card_no_separator(self):
        result = mask_text("Card: 4111111111111111", role="Analyst", rules_config=RULES)
        self.assertNotIn("4111111111111111", result)

    def test_phone_masked(self):
        result = mask_text("Phone: 555-867-5309", role="Analyst", rules_config=RULES)
        self.assertNotIn("555-867-5309", result)
        self.assertIn("***-***-****", result)

    def test_phone_dot_separator(self):
        result = mask_text("Phone: 555.867.5309", role="Analyst", rules_config=RULES)
        self.assertNotIn("555.867.5309", result)

    def test_admin_sees_all(self):
        text = "SSN: 123-45-6789, email: a@b.com"
        result = mask_text(text, role="Admin", rules_config=RULES)
        self.assertEqual(result, text)

    def test_multiple_pii_in_one_string(self):
        text = "SSN: 123-45-6789 and email: user@test.org"
        result = mask_text(text, role="Analyst", rules_config=RULES)
        self.assertNotIn("123-45-6789", result)
        self.assertNotIn("user@test.org", result)

    def test_name_masking_off_by_default(self):
        text = "Customer [NAME] has a balance"
        result = mask_text(text, role="Analyst", mask_names=False, rules_config=RULES)
        self.assertIn("[NAME]", result)

    def test_name_masking_on(self):
        text = "Customer [NAME] has a balance"
        result = mask_text(text, role="Analyst", mask_names=True, rules_config=RULES)
        self.assertNotIn("[NAME]", result)

    def test_no_false_positives(self):
        text = "The year 2026 and zip code 90210 should not be masked"
        result = mask_text(text, role="Analyst", rules_config=RULES)
        self.assertIn("2026", result)
        self.assertIn("90210", result)


class TestMaskCSV(unittest.TestCase):

    CSV_CONTENT = (
        "id,name,ssn,email,phone\n"
        "1,John Doe,123-45-6789,john@example.com,555-123-4567\n"
        "2,Jane Smith,987-65-4321,jane@example.com,555-987-6543\n"
    )

    def test_header_preserved(self):
        result = mask_csv(self.CSV_CONTENT, role="Analyst", rules_config=RULES)
        first_line = result.splitlines()[0]
        self.assertEqual(first_line, "id,name,ssn,email,phone")

    def test_ssn_masked_in_data_rows(self):
        result = mask_csv(self.CSV_CONTENT, role="Analyst", rules_config=RULES)
        self.assertNotIn("123-45-6789", result)
        self.assertNotIn("987-65-4321", result)

    def test_email_masked_in_data_rows(self):
        result = mask_csv(self.CSV_CONTENT, role="Analyst", rules_config=RULES)
        self.assertNotIn("john@example.com", result)
        self.assertNotIn("jane@example.com", result)

    def test_admin_passthrough(self):
        result = mask_csv(self.CSV_CONTENT, role="Admin", rules_config=RULES)
        self.assertIn("123-45-6789", result)
        self.assertIn("john@example.com", result)

    def test_row_count_preserved(self):
        result = mask_csv(self.CSV_CONTENT, role="Analyst", rules_config=RULES)
        self.assertEqual(len(result.splitlines()), len(self.CSV_CONTENT.splitlines()))


class TestMaskTxt(unittest.TestCase):

    def test_pii_masked(self):
        text = "Call me at 555-111-2222 or email me at foo@bar.com"
        result = mask_txt(text, role="Analyst", rules_config=RULES)
        self.assertNotIn("555-111-2222", result)
        self.assertNotIn("foo@bar.com", result)

    def test_non_pii_preserved(self):
        text = "The report was filed on 2026-03-12 in department 42."
        result = mask_txt(text, role="Analyst", rules_config=RULES)
        self.assertIn("2026-03-12", result)
        self.assertIn("department 42", result)


class TestMaskSQL(unittest.TestCase):

    SQL_CONTENT = (
        "INSERT INTO users VALUES (1, 'Alice', '111-22-3333', 'alice@corp.com', '555-000-1111');\n"
        "INSERT INTO users VALUES (2, 'Bob', '444-55-6666', 'bob@corp.com', '555-000-2222');\n"
    )

    def test_ssn_masked(self):
        result = mask_sql(self.SQL_CONTENT, role="Analyst", rules_config=RULES)
        self.assertNotIn("111-22-3333", result)
        self.assertNotIn("444-55-6666", result)

    def test_email_masked(self):
        result = mask_sql(self.SQL_CONTENT, role="Analyst", rules_config=RULES)
        self.assertNotIn("alice@corp.com", result)
        self.assertNotIn("bob@corp.com", result)

    def test_sql_keywords_preserved(self):
        result = mask_sql(self.SQL_CONTENT, role="Analyst", rules_config=RULES)
        self.assertIn("INSERT INTO users", result)

    def test_admin_passthrough(self):
        result = mask_sql(self.SQL_CONTENT, role="Admin", rules_config=RULES)
        self.assertIn("111-22-3333", result)


class TestMaskFile(unittest.TestCase):

    def _write_temp(self, suffix, content):
        fd, path = tempfile.mkstemp(suffix=suffix)
        os.close(fd)
        with open(path, "w", encoding="utf-8") as f:
            f.write(content)
        return path

    def test_csv_file(self):
        content = "id,ssn\n1,123-45-6789\n"
        path = self._write_temp(".csv", content)
        try:
            result = mask_file(path, role="Analyst", rules_config=RULES)
            self.assertNotIn("123-45-6789", result)
        finally:
            os.unlink(path)

    def test_txt_file(self):
        content = "Contact: 555-111-9999 or user@test.net\n"
        path = self._write_temp(".txt", content)
        try:
            result = mask_file(path, role="Analyst", rules_config=RULES)
            self.assertNotIn("555-111-9999", result)
            self.assertNotIn("user@test.net", result)
        finally:
            os.unlink(path)

    def test_sql_file(self):
        content = "INSERT INTO t VALUES ('abc@def.com', '222-33-4444');\n"
        path = self._write_temp(".sql", content)
        try:
            result = mask_file(path, role="Analyst", rules_config=RULES)
            self.assertNotIn("abc@def.com", result)
            self.assertNotIn("222-33-4444", result)
        finally:
            os.unlink(path)

    def test_output_written_to_disk(self):
        content = "SSN: 123-45-6789\n"
        path = self._write_temp(".txt", content)
        out_path = path + ".masked"
        try:
            mask_file(path, role="Analyst", output_path=out_path, rules_config=RULES)
            self.assertTrue(os.path.exists(out_path))
            with open(out_path, "r") as f:
                result = f.read()
            self.assertNotIn("123-45-6789", result)
        finally:
            os.unlink(path)
            if os.path.exists(out_path):
                os.unlink(out_path)

    def test_file_not_found_raises(self):
        with self.assertRaises(FileNotFoundError):
            mask_file("/nonexistent/path/file.csv")

    def test_unknown_extension_treated_as_text(self):
        # Unknown/no extension falls back to plain-text masking rather than raising.
        path = self._write_temp(".xyz", "Call 555-867-5309 for info")
        try:
            result = mask_file(path)
            self.assertIn("***-***-****", result)
        finally:
            os.unlink(path)

    def test_pdf_without_output_path_raises(self):
        path = self._write_temp(".pdf", "")
        try:
            with self.assertRaises(ValueError):
                mask_file(path, role="Analyst")
        finally:
            os.unlink(path)


class TestMaskPDF(unittest.TestCase):

    def _make_pdf_with_text(self, text: str) -> str:
        fd, path = tempfile.mkstemp(suffix=".pdf")
        os.close(fd)
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((50, 100), text, fontsize=11)
        doc.save(path)
        doc.close()
        return path

    def test_pdf_ssn_redacted(self):
        ssn = "123-45-6789"
        in_path = self._make_pdf_with_text(f"SSN: {ssn}")
        fd, out_path = tempfile.mkstemp(suffix=".pdf")
        os.close(fd)
        try:
            from core.masking import mask_pdf
            mask_pdf(in_path, out_path, role="Analyst", rules_config=RULES)
            doc = fitz.open(out_path)
            text = "".join(page.get_text() for page in doc)
            doc.close()
            self.assertNotIn(ssn, text)
        finally:
            os.unlink(in_path)
            os.unlink(out_path)

    def test_pdf_missing_input_raises(self):
        from core.masking import mask_pdf
        with self.assertRaises(FileNotFoundError):
            mask_pdf("/no/such/file.pdf", "/tmp/out.pdf", rules_config=RULES)


class TestMaskXLSX(unittest.TestCase):

    def _make_xlsx(self, ssn, email):
        from openpyxl import Workbook
        fd, path = tempfile.mkstemp(suffix=".xlsx")
        os.close(fd)
        wb = Workbook()
        ws = wb.active
        ws.append(["id", "ssn", "email"])
        ws.append([1, ssn, email])
        wb.save(path)
        return path

    def test_xlsx_pii_masked(self):
        from core.masking import mask_xlsx
        from openpyxl import load_workbook
        ssn = "123-45-6789"
        email = "user@example.com"
        in_path = self._make_xlsx(ssn, email)
        fd, out_path = tempfile.mkstemp(suffix=".xlsx")
        os.close(fd)
        try:
            mask_xlsx(in_path, out_path, role="Analyst", rules_config=RULES)
            wb = load_workbook(out_path)
            ws = wb.active
            vals = [str(cell.value or "") for row in ws.iter_rows(min_row=2) for cell in row]
            self.assertFalse(any(ssn in v for v in vals))
            self.assertFalse(any(email in v for v in vals))
        finally:
            os.unlink(in_path)
            os.unlink(out_path)

    def test_xlsx_header_preserved(self):
        from core.masking import mask_xlsx
        from openpyxl import load_workbook
        in_path = self._make_xlsx("111-22-3333", "a@b.com")
        fd, out_path = tempfile.mkstemp(suffix=".xlsx")
        os.close(fd)
        try:
            mask_xlsx(in_path, out_path, role="Analyst", rules_config=RULES)
            wb = load_workbook(out_path)
            ws = wb.active
            headers = [cell.value for cell in ws[1]]
            self.assertEqual(headers, ["id", "ssn", "email"])
        finally:
            os.unlink(in_path)
            os.unlink(out_path)

    def test_mask_file_xlsx_dispatch(self):
        in_path = self._make_xlsx("999-88-7777", "pii@test.com")
        fd, out_path = tempfile.mkstemp(suffix=".xlsx")
        os.close(fd)
        try:
            result = mask_file(in_path, role="Analyst", output_path=out_path, rules_config=RULES)
            self.assertIsNone(result)
            self.assertTrue(os.path.exists(out_path))
        finally:
            os.unlink(in_path)
            os.unlink(out_path)


class TestMaskDOCX(unittest.TestCase):

    def _make_docx(self, text):
        from docx import Document
        fd, path = tempfile.mkstemp(suffix=".docx")
        os.close(fd)
        doc = Document()
        doc.add_paragraph(text)
        doc.save(path)
        return path

    def test_docx_pii_masked(self):
        from core.masking import mask_docx
        from docx import Document
        ssn = "123-45-6789"
        in_path = self._make_docx(f"SSN: {ssn}, email: user@test.com")
        fd, out_path = tempfile.mkstemp(suffix=".docx")
        os.close(fd)
        try:
            mask_docx(in_path, out_path, role="Analyst", rules_config=RULES)
            doc = Document(out_path)
            full_text = " ".join(p.text for p in doc.paragraphs)
            self.assertNotIn(ssn, full_text)
            self.assertNotIn("user@test.com", full_text)
        finally:
            os.unlink(in_path)
            os.unlink(out_path)

    def test_mask_file_docx_dispatch(self):
        in_path = self._make_docx("Phone: 555-123-4567")
        fd, out_path = tempfile.mkstemp(suffix=".docx")
        os.close(fd)
        try:
            result = mask_file(in_path, role="Analyst", output_path=out_path, rules_config=RULES)
            self.assertIsNone(result)
            self.assertTrue(os.path.exists(out_path))
        finally:
            os.unlink(in_path)
            os.unlink(out_path)


class TestMaskODT(unittest.TestCase):

    def _make_odt(self, text):
        from odf.opendocument import OpenDocumentText
        from odf.text import P
        fd, path = tempfile.mkstemp(suffix=".odt")
        os.close(fd)
        odt = OpenDocumentText()
        odt.text.addElement(P(text=text))
        odt.save(path)
        return path

    def test_odt_pii_masked(self):
        from core.masking import mask_odt
        from odf.opendocument import load as odf_load
        from odf import teletype
        from odf.text import P
        ssn = "123-45-6789"
        in_path = self._make_odt(f"SSN: {ssn}")
        fd, out_path = tempfile.mkstemp(suffix=".odt")
        os.close(fd)
        try:
            mask_odt(in_path, out_path, role="Analyst", rules_config=RULES)
            doc = odf_load(out_path)
            text = " ".join(teletype.extractText(p) for p in doc.getElementsByType(P))
            self.assertNotIn(ssn, text)
        finally:
            os.unlink(in_path)
            os.unlink(out_path)


class TestMaskODS(unittest.TestCase):

    def _make_ods(self, text):
        from odf.opendocument import OpenDocumentSpreadsheet
        from odf.table import Table, TableRow, TableCell
        from odf.text import P
        fd, path = tempfile.mkstemp(suffix=".ods")
        os.close(fd)
        ods = OpenDocumentSpreadsheet()
        table = Table(name="Sheet1")
        row = TableRow()
        cell = TableCell()
        cell.addElement(P(text=text))
        row.addElement(cell)
        table.addElement(row)
        ods.spreadsheet.addElement(table)
        ods.save(path)
        return path

    def test_ods_pii_masked(self):
        from core.masking import mask_ods
        from odf.opendocument import load as odf_load
        from odf.table import TableCell
        from odf import teletype
        from odf.text import P
        ssn = "987-65-4321"
        in_path = self._make_ods(f"SSN: {ssn}")
        fd, out_path = tempfile.mkstemp(suffix=".ods")
        os.close(fd)
        try:
            mask_ods(in_path, out_path, role="Analyst", rules_config=RULES)
            doc = odf_load(out_path)
            text = " ".join(teletype.extractText(p) for p in doc.getElementsByType(P))
            self.assertNotIn(ssn, text)
        finally:
            os.unlink(in_path)
            os.unlink(out_path)


class TestMaskPPTX(unittest.TestCase):

    def _make_pptx(self, text):
        from pptx import Presentation
        fd, path = tempfile.mkstemp(suffix=".pptx")
        os.close(fd)
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Test"
        slide.placeholders[1].text = text
        prs.save(path)
        return path

    def test_pptx_pii_masked(self):
        from core.masking import mask_pptx
        from pptx import Presentation
        ssn = "111-22-3333"
        in_path = self._make_pptx(f"SSN: {ssn}, email: foo@bar.com")
        fd, out_path = tempfile.mkstemp(suffix=".pptx")
        os.close(fd)
        try:
            mask_pptx(in_path, out_path, role="Analyst", rules_config=RULES)
            prs = Presentation(out_path)
            all_text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        all_text.append(shape.text_frame.text)
            full = " ".join(all_text)
            self.assertNotIn(ssn, full)
            self.assertNotIn("foo@bar.com", full)
        finally:
            os.unlink(in_path)
            os.unlink(out_path)


class TestMaskFileGuards(unittest.TestCase):

    def _write_temp(self, suffix, content=b"dummy"):
        fd, path = tempfile.mkstemp(suffix=suffix)
        os.close(fd)
        with open(path, "wb") as f:
            f.write(content)
        return path

    def test_all_encrypt_only_extensions_raise(self):
        from core.masking import _ENCRYPT_ONLY_EXTS
        for ext in sorted(_ENCRYPT_ONLY_EXTS):
            path = self._write_temp(ext)
            try:
                with self.assertRaises(ValueError) as ctx:
                    mask_file(path, rules_config=RULES)
                self.assertIn(
                    "encrypt-only", str(ctx.exception),
                    msg=f"Expected 'encrypt-only' in error for {ext}"
                )
            finally:
                os.unlink(path)

    def test_doc_raises(self):
        path = self._write_temp(".doc")
        try:
            with self.assertRaises(ValueError) as ctx:
                mask_file(path, rules_config=RULES)
            self.assertIn(".docx", str(ctx.exception))
        finally:
            os.unlink(path)

    def test_ppt_raises(self):
        path = self._write_temp(".ppt")
        try:
            with self.assertRaises(ValueError) as ctx:
                mask_file(path, rules_config=RULES)
            self.assertIn(".pptx", str(ctx.exception))
        finally:
            os.unlink(path)

    def test_rtf_raises(self):
        path = self._write_temp(".rtf")
        try:
            with self.assertRaises(ValueError) as ctx:
                mask_file(path, rules_config=RULES)
            self.assertIn(".rtf", str(ctx.exception))
        finally:
            os.unlink(path)

    def test_tex_treated_as_text(self):
        fd, path = tempfile.mkstemp(suffix=".tex")
        os.close(fd)
        with open(path, "w") as f:
            f.write(r"Author SSN: 123-45-6789 \textbf{bold}")
        try:
            result = mask_file(path, role="Analyst", rules_config=RULES)
            self.assertNotIn("123-45-6789", result)
            self.assertIn(r"\textbf{bold}", result)
        finally:
            os.unlink(path)

    def test_bib_treated_as_text(self):
        fd, path = tempfile.mkstemp(suffix=".bib")
        os.close(fd)
        with open(path, "w") as f:
            f.write("@article{key, author={Jane user@test.com Doe}}")
        try:
            result = mask_file(path, role="Analyst", rules_config=RULES)
            self.assertNotIn("user@test.com", result)
        finally:
            os.unlink(path)

    def test_multiline_csv_quoted_field(self):
        content = 'id,note\n1,"line one\nSSN: 123-45-6789\nline three"\n'
        result = mask_csv(content, role="Analyst", rules_config=RULES)
        self.assertNotIn("123-45-6789", result)
        self.assertIn("line one", result)
        self.assertIn("line three", result)

    def test_malformed_csv_falls_back_to_text_masking(self):
        """A .csv file with NUL bytes (bad binary content) should fall back to
        plain-text masking instead of raising an exception."""
        fd, path = tempfile.mkstemp(suffix=".csv")
        os.close(fd)
        with open(path, "wb") as f:
            f.write(b"id,ssn\n1,123-45-6789\x00garbage\n")
        try:
            result = mask_file(path, role="Analyst", rules_config=RULES)
            self.assertIsNotNone(result)
            self.assertNotIn("123-45-6789", result)
        finally:
            os.unlink(path)


class TestMaskPDFTextPath(unittest.TestCase):

    def _make_text_pdf(self, text):
        fd, path = tempfile.mkstemp(suffix=".pdf")
        os.close(fd)
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((50, 100), text, fontsize=11)
        doc.save(path)
        doc.close()
        return path

    def test_text_path_ssn_redacted(self):
        from core.masking import mask_pdf
        ssn = "123-45-6789"
        in_path = self._make_text_pdf(f"SSN: {ssn}")
        fd, out_path = tempfile.mkstemp(suffix=".pdf")
        os.close(fd)
        try:
            mask_pdf(in_path, out_path, role="Analyst", rules_config=RULES)
            doc = fitz.open(out_path)
            text = "".join(page.get_text() for page in doc)
            doc.close()
            self.assertNotIn(ssn, text)
        finally:
            os.unlink(in_path)
            os.unlink(out_path)

    def test_mask_file_pdf_requires_output_path(self):
        in_path = self._make_text_pdf("SSN: 123-45-6789")
        try:
            with self.assertRaises(ValueError):
                mask_file(in_path, role="Analyst", rules_config=RULES)
        finally:
            os.unlink(in_path)


class TestMaskPDFOCRPath(unittest.TestCase):
    """Tests for the OCR (image-based) PDF masking path."""

    def _make_image_pdf(self, text):
        """Create a PDF with no text layer — content is a rasterised image."""
        from PIL import Image, ImageDraw, ImageFont
        fd, path = tempfile.mkstemp(suffix=".pdf")
        os.close(fd)
        img = Image.new("RGB", (800, 200), color=(255, 255, 255))
        draw = ImageDraw.Draw(img)
        draw.text((20, 80), text, fill=(0, 0, 0))
        img.save(path, format="PDF")
        return path

    def test_ocr_path_taken_for_image_pdf(self):
        from core.masking import _is_image_based_page
        in_path = self._make_image_pdf("SSN: 123-45-6789")
        try:
            doc = fitz.open(in_path)
            self.assertTrue(
                any(_is_image_based_page(p) for p in doc),
                "Expected at least one image-based page"
            )
            doc.close()
        finally:
            os.unlink(in_path)

    def test_ocr_path_output_is_raster_pdf(self):
        from core.masking import mask_pdf
        ssn = "123-45-6789"
        in_path = self._make_image_pdf(f"SSN: {ssn}")
        fd, out_path = tempfile.mkstemp(suffix=".pdf")
        os.close(fd)
        try:
            mask_pdf(in_path, out_path, role="Analyst", rules_config=RULES)
            self.assertTrue(os.path.exists(out_path))
            self.assertGreater(os.path.getsize(out_path), 0)
            # Output is raster-only — text layer should be empty
            doc = fitz.open(out_path)
            text = "".join(page.get_text() for page in doc)
            doc.close()
            self.assertEqual(text.strip(), "")
        finally:
            os.unlink(in_path)
            os.unlink(out_path)

    def test_ocr_path_ssn_region_blacked_out(self):
        """Run OCR on the masked output and confirm SSN pattern no longer found."""
        import pytesseract
        from core.masking import mask_pdf
        ssn = "123-45-6789"
        in_path = self._make_image_pdf(f"SSN: {ssn}")
        fd, out_path = tempfile.mkstemp(suffix=".pdf")
        os.close(fd)
        try:
            mask_pdf(in_path, out_path, role="Analyst", rules_config=RULES)
            doc = fitz.open(out_path)
            page = doc[0]
            pix = page.get_pixmap(dpi=300)
            from PIL import Image
            img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
            doc.close()
            ocr_text = pytesseract.image_to_string(img)
            self.assertNotIn(ssn, ocr_text)
        finally:
            os.unlink(in_path)
            os.unlink(out_path)


if __name__ == "__main__":
    unittest.main(verbosity=2)
