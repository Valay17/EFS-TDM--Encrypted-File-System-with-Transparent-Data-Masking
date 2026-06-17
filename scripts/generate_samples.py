"""
Generate synthetic test files containing PII using Faker.
Outputs: CSV, TXT, SQL dump, PDF, XLSX, DOCX, ODT, ODS, ODP, PPTX.

Usage (from project root EFS_TDM/):
    python scripts/generate_samples.py
    python scripts/generate_samples.py --count 20

Usage (from server_pkg/):
    python ../scripts/generate_samples.py
    python ../scripts/generate_samples.py --output-dir data/sample --count 20
"""

import argparse
import random
from pathlib import Path

import fitz  # PyMuPDF
from faker import Faker
from openpyxl import Workbook
from docx import Document
from odf.opendocument import OpenDocumentText, OpenDocumentSpreadsheet, OpenDocumentPresentation
from odf.text import P
from odf.table import Table, TableRow, TableCell, TableColumn
from odf.draw import Page as DrawPage, Frame, TextBox
from pptx import Presentation
from pptx.util import Inches, Pt

fake = Faker("en_US")
Faker.seed(42)
random.seed(42)


def _ssn():
    return f"{random.randint(100,999):03d}-{random.randint(10,99):02d}-{random.randint(1000,9999):04d}"


def _cc():
    groups = [f"{random.randint(1000,9999):04d}" for _ in range(4)]
    sep = random.choice(["-", " ", ""])
    return sep.join(groups)


def _phone():
    return f"{random.randint(200,999)}-{random.randint(200,999)}-{random.randint(1000,9999):04d}"


def generate_csv(output_path: str, count: int = 10):
    lines = ["id,name,ssn,email,phone,credit_card,salary"]
    for i in range(1, count + 1):
        row = ",".join([
            str(i),
            fake.name().replace(",", ""),
            _ssn(),
            fake.email(),
            _phone(),
            _cc(),
            str(random.randint(40000, 150000)),
        ])
        lines.append(row)
    Path(output_path).write_text("\n".join(lines) + "\n", encoding="utf-8")
    print(f"  CSV  -> {output_path}")


def generate_txt(output_path: str, count: int = 10):
    paragraphs = []
    for _ in range(count):
        name = fake.name()
        paragraphs.append(
            f"Customer {name} can be reached at {fake.email()} or {_phone()}. "
            f"Their SSN on file is {_ssn()} and their registered card ends in "
            f"{_cc()}. Account was opened on {fake.date()}."
        )
    Path(output_path).write_text("\n\n".join(paragraphs) + "\n", encoding="utf-8")
    print(f"  TXT  -> {output_path}")


def generate_sql(output_path: str, count: int = 10):
    lines = [
        "-- Synthetic employee dump",
        "CREATE TABLE IF NOT EXISTS employees (",
        "  id INTEGER PRIMARY KEY,",
        "  name TEXT,",
        "  ssn TEXT,",
        "  email TEXT,",
        "  phone TEXT,",
        "  credit_card TEXT,",
        "  salary INTEGER",
        ");",
        "",
        "BEGIN TRANSACTION;",
    ]
    for i in range(1, count + 1):
        name = fake.name().replace("'", "''")
        lines.append(
            f"INSERT INTO employees VALUES ("
            f"{i}, '{name}', '{_ssn()}', '{fake.email()}', "
            f"'{_phone()}', '{_cc()}', {random.randint(40000, 150000)});"
        )
    lines += ["COMMIT;", ""]
    Path(output_path).write_text("\n".join(lines), encoding="utf-8")
    print(f"  SQL  -> {output_path}")


def generate_xlsx(output_path: str, count: int = 10):
    wb = Workbook()
    ws = wb.active
    ws.title = "Employees"
    ws.append(["id", "name", "ssn", "email", "phone", "credit_card", "salary"])
    for i in range(1, count + 1):
        ws.append([
            i,
            fake.name(),
            _ssn(),
            fake.email(),
            _phone(),
            _cc(),
            random.randint(40000, 150000),
        ])
    wb.save(output_path)
    print(f"  XLSX -> {output_path}")


def generate_docx(output_path: str, count: int = 10):
    doc = Document()
    doc.add_heading("EFS-TDM Synthetic PII Report", 0)
    for i in range(1, count + 1):
        name = fake.name()
        doc.add_paragraph(
            f"{i}. {name} | SSN: {_ssn()} | Email: {fake.email()} | "
            f"Phone: {_phone()} | Card: {_cc()}"
        )
    doc.save(output_path)
    print(f"  DOCX -> {output_path}")


def generate_odt(output_path: str, count: int = 10):
    odt = OpenDocumentText()
    for i in range(1, count + 1):
        name = fake.name()
        p = P(text=(
            f"{i}. {name} | SSN: {_ssn()} | Email: {fake.email()} | "
            f"Phone: {_phone()} | Card: {_cc()}"
        ))
        odt.text.addElement(p)
    odt.save(output_path)
    print(f"  ODT  -> {output_path}")


def generate_ods(output_path: str, count: int = 10):
    ods = OpenDocumentSpreadsheet()
    table = Table(name="Employees")
    header_row = TableRow()
    for col in ["id", "name", "ssn", "email", "phone", "credit_card", "salary"]:
        cell = TableCell()
        cell.addElement(P(text=col))
        header_row.addElement(cell)
    table.addElement(header_row)
    for i in range(1, count + 1):
        row = TableRow()
        for val in [str(i), fake.name(), _ssn(), fake.email(), _phone(), _cc(), str(random.randint(40000, 150000))]:
            cell = TableCell()
            cell.addElement(P(text=val))
            row.addElement(cell)
        table.addElement(row)
    ods.spreadsheet.addElement(table)
    ods.save(output_path)
    print(f"  ODS  -> {output_path}")


def generate_odp(output_path: str, count: int = 10):
    odp = OpenDocumentPresentation()
    for i in range(1, count + 1):
        name = fake.name()
        page = DrawPage(stylename="", masterpagename="")
        frame = Frame(width="20cm", height="2cm", x="2cm", y="2cm")
        tb = TextBox()
        tb.addElement(P(text=(
            f"{i}. {name} | SSN: {_ssn()} | Email: {fake.email()} | "
            f"Phone: {_phone()} | Card: {_cc()}"
        )))
        frame.addElement(tb)
        page.addElement(frame)
        odp.presentation.addElement(page)
    odp.save(output_path)
    print(f"  ODP  -> {output_path}")


def generate_pptx(output_path: str, count: int = 10):
    prs = Presentation()
    layout = prs.slide_layouts[1]
    for i in range(1, count + 1):
        name = fake.name()
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = f"Record {i}"
        body = slide.placeholders[1]
        body.text = (
            f"{name} | SSN: {_ssn()} | Email: {fake.email()} | "
            f"Phone: {_phone()} | Card: {_cc()}"
        )
    prs.save(output_path)
    print(f"  PPTX -> {output_path}")


def generate_image_pdf(output_path: str, count: int = 10):
    """Image-based PDF (no text layer) — triggers OCR masking path."""
    from PIL import Image, ImageDraw
    pages = []
    for i in range(1, min(count, 3) + 1):
        img = Image.new("RGB", (1200, 200), color=(255, 255, 255))
        draw = ImageDraw.Draw(img)
        name = fake.name()
        line = (
            f"{i}. {name} | SSN: {_ssn()} | Email: {fake.email()} | "
            f"Phone: {_phone()} | Card: {_cc()}"
        )
        draw.text((20, 80), line, fill=(0, 0, 0))
        pages.append(img)
    pages[0].save(
        output_path,
        save_all=True,
        append_images=pages[1:],
        format="PDF",
    )
    print(f"  PDF (image) -> {output_path}")


def generate_pdf(output_path: str, count: int = 10):
    doc = fitz.open()
    page = doc.new_page()
    y = 50
    page.insert_text((50, y), "EFS-TDM Synthetic PII Report", fontsize=14)
    y += 30


    for i in range(1, count + 1):
        name = fake.name()
        line = (
            f"{i}. {name} | SSN: {_ssn()} | Email: {fake.email()} | "
            f"Phone: {_phone()} | Card: {_cc()}"
        )
        page.insert_text((50, y), line, fontsize=9)
        y += 16
        if y > 780:
            page = doc.new_page()
            y = 50

    doc.save(output_path)
    doc.close()
    print(f"  PDF  -> {output_path}")


_SCRIPT_DIR = Path(__file__).parent
_DEFAULT_OUT = _SCRIPT_DIR.parent / "server_pkg" / "data" / "sample"


def main():
    parser = argparse.ArgumentParser(description="Generate synthetic PII test files")
    parser.add_argument(
        "--output-dir", default=str(_DEFAULT_OUT),
        help="Directory to write sample files (default: server_pkg/data/sample)"
    )
    parser.add_argument(
        "--count", type=int, default=10, help="Number of records per file"
    )
    args = parser.parse_args()

    out = Path(args.output_dir)
    out.mkdir(parents=True, exist_ok=True)

    print(f"Generating {args.count} records into {out}/")
    generate_csv(str(out / "employees.csv"), args.count)
    generate_txt(str(out / "notes.txt"), args.count)
    generate_sql(str(out / "employees_dump.sql"), args.count)
    generate_pdf(str(out / "report.pdf"), args.count)
    generate_xlsx(str(out / "employees.xlsx"), args.count)
    generate_docx(str(out / "report.docx"), args.count)
    generate_odt(str(out / "report.odt"), args.count)
    generate_ods(str(out / "employees.ods"), args.count)
    generate_odp(str(out / "report.odp"), args.count)
    generate_pptx(str(out / "report.pptx"), args.count)
    generate_image_pdf(str(out / "report_scanned.pdf"), args.count)
    print("Done.")


if __name__ == "__main__":
    main()
